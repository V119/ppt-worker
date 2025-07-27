import re
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Tuple

from pptx import Presentation
from pptx.text.text import _Paragraph

from src import template_dir, output_dir


class TemplateRenderer:
    """处理Jinja2-like模板渲染和位置计算。"""
    
    PLACEHOLDER_PATTERN = re.compile(r'\{\{(.+?)}}')
    
    def __init__(self, context: Dict[str, any]):
        self.context = context
    
    def cal_render_position(self, all_text: str) -> Tuple[str, List[Dict[str, any]]]:
        """计算渲染后的文本和占位符位置信息。"""
        matches = self.PLACEHOLDER_PATTERN.finditer(all_text)
        placeholder_list: List[Dict[str, any]] = []
        rendered_text = ''
        last_end = 0
        
        for match in matches:
            rendered_text += all_text[last_end:match.start()]
            placeholder_text = match.group(1).strip()
            rendered_value = str(self.context.get(placeholder_text, ''))
            original_start = match.start()
            original_length = match.end() - match.start()
            rendered_start = len(rendered_text)
            rendered_length = len(rendered_value)
            placeholder_list.append({
                'original_start': original_start,
                'original_length': original_length,
                'rendered_start': rendered_start,
                'rendered_length': rendered_length,
                'placeholder': placeholder_text,
                'rendered_value': rendered_value
            })
            rendered_text += rendered_value
            last_end = match.end()
        
        rendered_text += all_text[last_end:]
        return rendered_text, placeholder_list

class PptProcessor:
    """处理PPT文件的加载、渲染和保存。"""
    
    def __init__(self, template_path: Path, output_dir: Path, renderer: TemplateRenderer):
        self.template_path = template_path
        self.output_dir = output_dir
        self.renderer = renderer
        self.prs = Presentation(str(self.template_path))
    
    def replace_paragraph_runs(self, paragraph: _Paragraph) -> None:
        """替换段落中的run文本，处理跨run占位符。"""
        text_position_list: List[Dict[str, any]] = []
        start_index = 0
        for run in paragraph.runs:
            text = run.text
            text_position_list.append({
                'start': start_index,
                'len': len(text),
                'run': run
            })
            start_index += len(text)
        
        all_text = ''.join(item['run'].text for item in text_position_list)
        rendered_text, placeholder_list = self.renderer.cal_render_position(all_text)
        
        if not placeholder_list:
            return
        
        new_texts = ['' for _ in text_position_list]
        current_pos = 0
        
        for ph in placeholder_list:
            orig_start = ph['original_start']
            self._copy_non_placeholder(new_texts, text_position_list, all_text, current_pos, orig_start)
            current_pos = orig_start + ph['original_length']
            self._allocate_rendered_value(new_texts, text_position_list, ph, orig_start, current_pos)
        
        self._copy_non_placeholder(new_texts, text_position_list, all_text, current_pos, len(all_text))
        
        for i, new_text in enumerate(new_texts):
            text_position_list[i]['run'].text = new_text
    
    @staticmethod
    def _copy_non_placeholder(new_texts: List[str], text_position_list: List[Dict], all_text: str, start: int, end: int) -> None:
        """复制非占位符文本到新文本列表。"""
        current = start
        while current < end:
            for idx, pos in enumerate(text_position_list):
                run_start = pos['start']
                run_end = run_start + pos['len']
                if run_end > current >= run_start:
                    rel_pos = current - run_start
                    remaining = min(end - current, run_end - current)
                    new_texts[idx] += all_text[current:current + remaining]
                    current += remaining
                    break
    
    @staticmethod
    def _allocate_rendered_value(new_texts: List[str], text_position_list: List[Dict], ph: Dict, orig_start: int, orig_end: int) -> None:
        """按比例分配渲染值到覆盖的run。"""
        covered_runs = []
        total_covered_len = 0
        for idx, pos in enumerate(text_position_list):
            run_start = pos['start']
            run_end = run_start + pos['len']
            if run_end > orig_start and run_start < orig_end:
                part_start = max(run_start, orig_start)
                part_end = min(run_end, orig_end)
                part_len = part_end - part_start
                covered_runs.append((idx, part_len))
                total_covered_len += part_len
        
        if total_covered_len == 0:
            return
        
        rend_value = ph['rendered_value']
        rend_len = ph['rendered_length']
        allocated = 0
        for i, (idx, part_len) in enumerate(covered_runs):
            ratio = part_len / total_covered_len
            alloc_len = round(ratio * rend_len) if i < len(covered_runs) - 1 else rend_len - allocated
            new_texts[idx] += rend_value[allocated:allocated + alloc_len]
            allocated += alloc_len
    
    def process(self) -> Path:
        """处理整个PPT并保存输出。"""
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    if not paragraph.runs:
                        continue
                    self.replace_paragraph_runs(paragraph)
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = f'output_{timestamp}.pptx'
        output_path = self.output_dir / output_filename
        self.prs.save(str(output_path))
        print(f'PPT已成功生成并保存到: {output_path}')
        return output_path

def main():
    context = {'sales': 980.0, 'growth': 18.98}
    template_file_path = template_dir / 'template.pptx'
    renderer = TemplateRenderer(context)
    processor = PptProcessor(template_file_path, output_dir, renderer)
    processor.process()

if __name__ == '__main__':
    main()

    

